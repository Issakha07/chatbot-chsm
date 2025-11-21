# -*- coding: utf-8 -*-
"""
Module de traitement universel de documents
Supporte : PDF, Word, Excel, TXT, PowerPoint, Images
"""

import os
import logging
from pathlib import Path
from typing import List, Dict, Optional
import fitz  # PyMuPDF
from docx import Document
import openpyxl
import pandas as pd
from pptx import Presentation
from PIL import Image
import pdfplumber

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """
    Processeur universel de documents pour RAG
    Extrait le texte de multiples formats de fichiers
    """
    
    SUPPORTED_EXTENSIONS = {
        '.pdf': 'PDF',
        '.docx': 'Word',
        '.doc': 'Word (ancien)',
        '.xlsx': 'Excel',
        '.xls': 'Excel (ancien)',
        '.csv': 'CSV',
        '.txt': 'Texte',
        '.pptx': 'PowerPoint',
        '.png': 'Image PNG',
        '.jpg': 'Image JPEG',
        '.jpeg': 'Image JPEG'
    }
    
    def __init__(self):
        self.stats = {
            'total_files': 0,
            'success': 0,
            'errors': 0,
            'by_type': {}
        }
    
    def is_supported(self, file_path: str) -> bool:
        """Vérifie si le format de fichier est supporté"""
        ext = Path(file_path).suffix.lower()
        return ext in self.SUPPORTED_EXTENSIONS
    
    def get_file_type(self, file_path: str) -> Optional[str]:
        """Retourne le type de fichier"""
        ext = Path(file_path).suffix.lower()
        return self.SUPPORTED_EXTENSIONS.get(ext)
    
    # ==========================================
    # EXTRACTION PDF
    # ==========================================
    
    def extract_from_pdf_pymupdf(self, pdf_path: str) -> str:
        """
        Extraction PDF avec PyMuPDF (recommandé)
        Meilleur pour PDFs complexes avec images/tableaux
        """
        try:
            doc = fitz.open(pdf_path)
            text = ""
            
            for page_num, page in enumerate(doc, 1):
                # Extraire texte
                page_text = page.get_text()
                text += f"\n--- Page {page_num} ---\n{page_text}"
                
                # Optionnel : extraire texte des images (OCR)
                # Nécessite pytesseract pour OCR complet
            
            doc.close()
            return text.strip()
        
        except Exception as e:
            logger.error(f"Erreur PyMuPDF avec {pdf_path}: {e}")
            return ""
    
    def extract_from_pdf_pdfplumber(self, pdf_path: str) -> str:
        """
        Extraction PDF avec pdfplumber (alternatif)
        Meilleur pour tableaux structurés
        """
        try:
            text = ""
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text() or ""
                    
                    # Extraire tableaux
                    tables = page.extract_tables()
                    if tables:
                        page_text += "\n\n[TABLEAUX DÉTECTÉS]\n"
                        for table in tables:
                            page_text += self._format_table(table) + "\n"
                    
                    text += f"\n--- Page {page_num} ---\n{page_text}"
            
            return text.strip()
        
        except Exception as e:
            logger.error(f"Erreur pdfplumber avec {pdf_path}: {e}")
            return ""
    
    def _format_table(self, table: List[List]) -> str:
        """Formate un tableau en texte lisible"""
        if not table:
            return ""
        
        formatted = []
        for row in table:
            row_text = " | ".join([str(cell) if cell else "" for cell in row])
            formatted.append(row_text)
        
        return "\n".join(formatted)
    
    # ==========================================
    # EXTRACTION WORD
    # ==========================================
    
    def extract_from_word(self, word_path: str) -> str:
        """Extraction de texte depuis Word (.docx)"""
        try:
            doc = Document(word_path)
            text = []
            
            # Paragraphes
            for para in doc.paragraphs:
                if para.text.strip():
                    text.append(para.text)
            
            # Tableaux
            for table in doc.tables:
                text.append("\n[TABLEAU]")
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    text.append(row_text)
            
            return "\n\n".join(text)
        
        except Exception as e:
            logger.error(f"Erreur Word avec {word_path}: {e}")
            return ""
    
    # ==========================================
    # EXTRACTION EXCEL
    # ==========================================
    
    def extract_from_excel(self, excel_path: str) -> str:
        """Extraction de texte depuis Excel (.xlsx, .xls)"""
        try:
            # Utiliser pandas pour flexibilité
            excel_file = pd.ExcelFile(excel_path)
            text = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                text.append(f"\n=== Feuille: {sheet_name} ===\n")
                
                # Convertir DataFrame en texte structuré
                # Limiter à 1000 lignes pour éviter surcharge
                if len(df) > 1000:
                    text.append(f"[ATTENTION: Feuille tronquée - {len(df)} lignes -> affichage de 1000]")
                    df = df.head(1000)
                
                # Format texte lisible
                text.append(df.to_string(index=False))
            
            return "\n\n".join(text)
        
        except Exception as e:
            logger.error(f"Erreur Excel avec {excel_path}: {e}")
            return ""
    
    # ==========================================
    # EXTRACTION CSV
    # ==========================================
    
    def extract_from_csv(self, csv_path: str) -> str:
        """Extraction depuis CSV"""
        try:
            df = pd.read_csv(csv_path, encoding='utf-8', nrows=1000)
            return f"=== Fichier CSV ===\n{df.to_string(index=False)}"
        
        except Exception as e:
            # Essayer avec encoding différent
            try:
                df = pd.read_csv(csv_path, encoding='latin-1', nrows=1000)
                return f"=== Fichier CSV ===\n{df.to_string(index=False)}"
            except:
                logger.error(f"Erreur CSV avec {csv_path}: {e}")
                return ""
    
    # ==========================================
    # EXTRACTION POWERPOINT
    # ==========================================
    
    def extract_from_powerpoint(self, pptx_path: str) -> str:
        """Extraction depuis PowerPoint (.pptx)"""
        try:
            prs = Presentation(pptx_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"\n=== Slide {slide_num} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            
            return "\n\n".join(text)
        
        except Exception as e:
            logger.error(f"Erreur PowerPoint avec {pptx_path}: {e}")
            return ""
    
    # ==========================================
    # EXTRACTION TEXTE SIMPLE
    # ==========================================
    
    def extract_from_txt(self, txt_path: str) -> str:
        """Extraction depuis fichier texte"""
        try:
            with open(txt_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Essayer avec encoding différent
            try:
                with open(txt_path, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Erreur TXT avec {txt_path}: {e}")
                return ""
        except Exception as e:
            logger.error(f"Erreur TXT avec {txt_path}: {e}")
            return ""
    
    # ==========================================
    # EXTRACTION IMAGE (OCR optionnel)
    # ==========================================
    
    def extract_from_image(self, image_path: str) -> str:
        """
        Extraction depuis image
        NOTE: OCR nécessite pytesseract (non inclus par défaut)
        """
        # Pour l'instant, retourne juste les métadonnées
        try:
            img = Image.open(image_path)
            metadata = f"Image: {img.format} | Taille: {img.size} | Mode: {img.mode}"
            
            # TODO: Ajouter OCR avec pytesseract si nécessaire
            # from pytesseract import image_to_string
            # text = image_to_string(img, lang='fra')
            
            return f"[IMAGE]\n{metadata}\n[OCR non activé - installez pytesseract pour extraction de texte]"
        
        except Exception as e:
            logger.error(f"Erreur Image avec {image_path}: {e}")
            return ""
    
    # ==========================================
    # MÉTHODE PRINCIPALE
    # ==========================================
    
    def extract_text(self, file_path: str, use_fallback: bool = True) -> Dict[str, str]:
        """
        Méthode principale d'extraction de texte
        
        Args:
            file_path: Chemin du fichier
            use_fallback: Si True, essaie méthode alternative en cas d'échec
        
        Returns:
            Dict avec 'text', 'file_type', 'success', 'error'
        """
        file_path = str(file_path)
        ext = Path(file_path).suffix.lower()
        file_type = self.get_file_type(file_path)
        
        result = {
            'file_path': file_path,
            'file_type': file_type,
            'text': '',
            'success': False,
            'error': None
        }
        
        if not self.is_supported(file_path):
            result['error'] = f"Format non supporté: {ext}"
            logger.warning(result['error'])
            return result
        
        try:
            # Routage selon type de fichier
            if ext == '.pdf':
                # Essayer PyMuPDF d'abord (plus rapide)
                text = self.extract_from_pdf_pymupdf(file_path)
                
                # Fallback vers pdfplumber si échec
                if use_fallback and not text.strip():
                    logger.info(f"Fallback vers pdfplumber pour {file_path}")
                    text = self.extract_from_pdf_pdfplumber(file_path)
            
            elif ext == '.docx':
                text = self.extract_from_word(file_path)
            
            elif ext in ['.xlsx', '.xls']:
                text = self.extract_from_excel(file_path)
            
            elif ext == '.csv':
                text = self.extract_from_csv(file_path)
            
            elif ext == '.pptx':
                text = self.extract_from_powerpoint(file_path)
            
            elif ext == '.txt':
                text = self.extract_from_txt(file_path)
            
            elif ext in ['.png', '.jpg', '.jpeg']:
                text = self.extract_from_image(file_path)
            
            else:
                result['error'] = f"Extracteur non implémenté pour {ext}"
                return result
            
            # Vérifier résultat
            if text and text.strip():
                result['text'] = text.strip()
                result['success'] = True
                self.stats['success'] += 1
                logger.info(f"✅ Extraction réussie: {Path(file_path).name} ({len(text)} caractères)")
            else:
                result['error'] = "Extraction vide"
                self.stats['errors'] += 1
                logger.warning(f"⚠️ Extraction vide: {Path(file_path).name}")
        
        except Exception as e:
            result['error'] = str(e)
            self.stats['errors'] += 1
            logger.error(f"❌ Erreur extraction {Path(file_path).name}: {e}")
        
        finally:
            self.stats['total_files'] += 1
            type_key = file_type or 'Unknown'
            self.stats['by_type'][type_key] = self.stats['by_type'].get(type_key, 0) + 1
        
        return result
    
    # ==========================================
    # TRAITEMENT BATCH
    # ==========================================
    
    def process_directory(self, directory: str, recursive: bool = False) -> List[Dict]:
        """
        Traite tous les fichiers supportés dans un dossier
        
        Args:
            directory: Chemin du dossier
            recursive: Si True, parcourt sous-dossiers
        
        Returns:
            Liste de résultats d'extraction
        """
        directory = Path(directory)
        
        if not directory.exists():
            logger.error(f"Dossier introuvable: {directory}")
            return []
        
        # Trouver tous les fichiers supportés
        pattern = '**/*' if recursive else '*'
        all_files = directory.glob(pattern)
        
        supported_files = [
            f for f in all_files 
            if f.is_file() and self.is_supported(str(f))
        ]
        
        logger.info(f"📂 Traitement de {len(supported_files)} fichiers dans {directory}")
        
        results = []
        for file_path in supported_files:
            result = self.extract_text(str(file_path))
            results.append(result)
        
        # Afficher stats
        logger.info(f"""
        ╔══════════════════════════════════════╗
        ║   STATISTIQUES D'EXTRACTION          ║
        ╠══════════════════════════════════════╣
        ║ Total fichiers: {self.stats['total_files']:3d}                 ║
        ║ Succès:         {self.stats['success']:3d}                 ║
        ║ Erreurs:        {self.stats['errors']:3d}                 ║
        ╚══════════════════════════════════════╝
        """)
        
        for doc_type, count in self.stats['by_type'].items():
            logger.info(f"  • {doc_type}: {count}")
        
        return results
    
    def get_stats(self) -> Dict:
        """Retourne les statistiques d'extraction"""
        return self.stats.copy()


# ==========================================
# FONCTIONS UTILITAIRES
# ==========================================

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """
    Découpe le texte en chunks avec overlap
    
    Args:
        text: Texte à découper
        chunk_size: Taille en mots
        overlap: Nombre de mots de chevauchement
    
    Returns:
        Liste de chunks
    """
    words = text.split()
    chunks = []
    
    for i in range(0, len(words), chunk_size - overlap):
        chunk = ' '.join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
    
    return chunks


# ==========================================
# EXEMPLE D'UTILISATION
# ==========================================

if __name__ == "__main__":
    # Configuration logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s [%(levelname)s] %(message)s'
    )
    
    # Test du processeur
    processor = DocumentProcessor()
    
    # Traiter un dossier
    results = processor.process_directory("../documents")
    
    # Afficher résultats
    for result in results:
        if result['success']:
            print(f"\n✅ {Path(result['file_path']).name}")
            print(f"Type: {result['file_type']}")
            print(f"Texte extrait: {len(result['text'])} caractères")
            print(f"Aperçu: {result['text'][:200]}...")
        else:
            print(f"\n❌ {Path(result['file_path']).name}")
            print(f"Erreur: {result['error']}")
